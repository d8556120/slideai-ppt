"""
PPT Generation Module
Creates professional PowerPoint files using python-pptx with multiple template themes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
import uuid


# ── Theme definitions ──────────────────────────────────────────────────────────

THEMES = {
    "business": {
        "name": "Business Professional",
        "bg_primary": RGBColor(0x0F, 0x2B, 0x46),       # Dark navy
        "bg_secondary": RGBColor(0x14, 0x38, 0x5C),      # Slightly lighter navy
        "accent": RGBColor(0x00, 0x96, 0xD6),            # Bright blue
        "accent2": RGBColor(0x00, 0xC9, 0xA7),           # Teal
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0xE0, 0xE6, 0xED),
        "bullet_color": RGBColor(0x00, 0x96, 0xD6),
        "subtitle_color": RGBColor(0x8B, 0xA5, 0xBE),
        "slide_num_color": RGBColor(0x5A, 0x7A, 0x96),
        "divider_color": RGBColor(0x00, 0x96, 0xD6),
    },
    "creative": {
        "name": "Creative Modern",
        "bg_primary": RGBColor(0x1A, 0x0A, 0x2E),        # Deep purple
        "bg_secondary": RGBColor(0x22, 0x10, 0x3A),
        "accent": RGBColor(0xA8, 0x55, 0xF7),            # Purple
        "accent2": RGBColor(0xF4, 0x72, 0xB6),            # Pink
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0xE2, 0xDA, 0xEF),
        "bullet_color": RGBColor(0xA8, 0x55, 0xF7),
        "subtitle_color": RGBColor(0xB0, 0x97, 0xCC),
        "slide_num_color": RGBColor(0x6B, 0x5A, 0x8A),
        "divider_color": RGBColor(0xA8, 0x55, 0xF7),
    },
    "minimal": {
        "name": "Minimal Clean",
        "bg_primary": RGBColor(0xFA, 0xFA, 0xFA),        # Off-white
        "bg_secondary": RGBColor(0xF0, 0xF0, 0xF0),
        "accent": RGBColor(0x18, 0x18, 0x18),             # Near-black
        "accent2": RGBColor(0x66, 0x66, 0x66),
        "title_color": RGBColor(0x11, 0x11, 0x11),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "bullet_color": RGBColor(0x11, 0x11, 0x11),
        "subtitle_color": RGBColor(0x77, 0x77, 0x77),
        "slide_num_color": RGBColor(0xAA, 0xAA, 0xAA),
        "divider_color": RGBColor(0x22, 0x22, 0x22),
    }
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, left, top, width, height, color: RGBColor, alpha=None):
    """Add a colored rectangle shape as a background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border
    if alpha is not None:
        solid_fill = shape.fill._fill
        a_elem = solid_fill.find(qn('a:solidFill'))
        if a_elem is not None:
            srgb = a_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_page_number(slide, slide_num: int, total: int, theme: dict):
    """Add page number in the bottom right corner."""
    _add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                  f"{slide_num} / {total}", 10, theme["slide_num_color"],
                  alignment=PP_ALIGN.RIGHT, font_name="Segoe UI Light")


def _add_left_accent_bar(slide, theme: dict):
    """Add a thin accent bar on the left side of content slides."""
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_HEIGHT, theme["accent"])


def _add_business_header(slide, theme: dict):
    """Add a subtle header bar for business template."""
    # Header background bar
    _add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.5), theme["bg_secondary"])
    # Thin accent line under header
    _add_shape_bg(slide, Inches(0), Inches(0.5), SLIDE_WIDTH, Inches(0.03), theme["accent"])
    # Branding text
    _add_text_box(slide, Inches(0.5), Inches(0.05), Inches(3), Inches(0.4),
                  "SlideAI", 9, theme["slide_num_color"],
                  bold=True, font_name="Segoe UI Semibold")


def _add_creative_decorations(slide, theme: dict):
    """Add gradient-like colored shapes as background decoration for creative template."""
    # Large circle top-right (semi-transparent)
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = theme["accent"]
    circle1.line.fill.background()
    try:
        sf = circle1.fill._fill
        a_elem = sf.find(qn('a:solidFill'))
        if a_elem is not None:
            srgb = a_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '8000')
    except Exception:
        pass

    # Small circle bottom-left (semi-transparent)
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.5), Inches(3.5), Inches(3.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = theme["accent2"]
    circle2.line.fill.background()
    try:
        sf = circle2.fill._fill
        a_elem = sf.find(qn('a:solidFill'))
        if a_elem is not None:
            srgb = a_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '6000')
    except Exception:
        pass


def _add_minimal_dividers(slide, theme: dict):
    """Add elegant thin lines as dividers for minimal template."""
    # Top thin line
    _add_shape_bg(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.01), theme["divider_color"], alpha=30)
    # Bottom thin line
    _add_shape_bg(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.01), theme["divider_color"], alpha=30)


def _apply_template_decorations(slide, template: str, theme: dict):
    """Apply template-specific decorations to a content slide."""
    if template == "business":
        _add_business_header(slide, theme)
    elif template == "creative":
        _add_creative_decorations(slide, theme)
    elif template == "minimal":
        _add_minimal_dividers(slide, theme)


def _create_title_slide(prs, content: dict, theme: dict, template: str = "business"):
    """Create the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, theme["bg_primary"])

    # Left accent bar
    _add_shape_bg(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, theme["accent"])

    # Decorative accent circle (top-right)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = theme["accent"]
    circle.line.fill.background()
    try:
        sf = circle.fill._fill
        a_elem = sf.find(qn('a:solidFill'))
        if a_elem is not None:
            srgb = a_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '15000')
    except Exception:
        pass

    # Creative template: extra decoration
    if template == "creative":
        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = theme["accent2"]
        circle2.line.fill.background()
        try:
            sf = circle2.fill._fill
            a_elem = sf.find(qn('a:solidFill'))
            if a_elem is not None:
                srgb = a_elem.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_elem.set('val', '10000')
        except Exception:
            pass

    # Title - large centered text
    title_text = content.get("title", "Presentation")
    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(2.2),
                  title_text, 44, theme["title_color"],
                  bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    # Decorative line under title
    line_width = min(Inches(4), Inches(len(title_text) * 0.08))
    line_left = Inches(6.666) - line_width / 2  # Center the line
    _add_shape_bg(slide, line_left, Inches(4.1), line_width, Inches(0.06), theme["accent"])

    # Subtitle - centered below line
    _add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.9),
                  content.get("subtitle", ""), 22, theme["subtitle_color"],
                  alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

    # Author - centered at bottom
    _add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
                  content.get("author", ""), 14, theme["slide_num_color"],
                  alignment=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Minimal template: add elegant horizontal rules
    if template == "minimal":
        _add_shape_bg(slide, Inches(2), Inches(1.5), Inches(9.3), Inches(0.01), theme["divider_color"], alpha=20)
        _add_shape_bg(slide, Inches(2), Inches(6.5), Inches(9.3), Inches(0.01), theme["divider_color"], alpha=20)


def _create_content_slide(prs, slide_data: dict, slide_num: int, total: int,
                          theme: dict, template: str = "business"):
    """Create a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["bg_primary"])

    # Apply template-specific decorations (drawn first so content is on top)
    _apply_template_decorations(slide, template, theme)

    # Left accent bar for all templates
    _add_left_accent_bar(slide, theme)

    # Top accent bar
    _add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05), theme["accent"])

    # Content offset depends on template
    content_top_offset = 0.0
    if template == "business":
        content_top_offset = 0.55  # Account for header bar

    # Slide title
    _add_text_box(slide, Inches(0.8), Inches(0.6 + content_top_offset), Inches(11), Inches(1),
                  slide_data.get("title", ""), 36, theme["title_color"],
                  bold=True, font_name="Segoe UI Semibold")

    # Divider under title
    _add_shape_bg(slide, Inches(0.8), Inches(1.55 + content_top_offset),
                  Inches(1.5), Inches(0.04), theme["accent"])

    # Bullet points
    bullets = slide_data.get("bullets", [])
    bullet_top = 2.0 + content_top_offset
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(bullet_top), Inches(11), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_before = Pt(14)
        p.space_after = Pt(10)

        # Check if this is a sub-point (starts with "- " or "  ")
        is_sub = bullet.startswith("- ") or bullet.startswith("  ")
        if is_sub:
            bullet = bullet.lstrip("- ")
            p.level = 1

            # Sub-point marker
            run_marker = p.add_run()
            run_marker.text = "    \u25E6  "  # Open circle for sub-points with indentation
            run_marker.font.size = Pt(10)
            run_marker.font.color.rgb = theme["bullet_color"]
            run_marker.font.name = "Segoe UI"

            # Sub-point text (slightly smaller)
            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(16)
            run_text.font.color.rgb = theme["text_color"]
            run_text.font.name = "Segoe UI"
        else:
            p.level = 0

            # Bullet marker
            run_marker = p.add_run()
            run_marker.text = "\u25CF  "  # Filled circle
            run_marker.font.size = Pt(10)
            run_marker.font.color.rgb = theme["bullet_color"]
            run_marker.font.name = "Segoe UI"

            # Bullet text
            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(18)
            run_text.font.color.rgb = theme["text_color"]
            run_text.font.name = "Segoe UI"

    # Speaker notes
    notes = slide_data.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
        for paragraph in notes_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

    # Page number (bottom right)
    total_with_title_closing = total + 2
    _add_page_number(slide, slide_num + 1, total_with_title_closing, theme)

    # Bottom decorative bar
    _add_shape_bg(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), theme["bg_secondary"])


def _create_closing_slide(prs, closing: dict, theme: dict, template: str = "business",
                          total_slides: int = 0):
    """Create the closing/thank-you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["bg_primary"])

    # Large accent circle center
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.5), Inches(0.5), Inches(4.3), Inches(4.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = theme["accent"]
    circle.line.fill.background()
    try:
        sf = circle.fill._fill
        a_elem = sf.find(qn('a:solidFill'))
        if a_elem is not None:
            srgb = a_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '10000')
    except Exception:
        pass

    # Creative template: extra decorations
    if template == "creative":
        _add_creative_decorations(slide, theme)

    # Minimal template: elegant dividers
    if template == "minimal":
        _add_shape_bg(slide, Inches(3), Inches(1.5), Inches(7.3), Inches(0.01), theme["divider_color"], alpha=20)
        _add_shape_bg(slide, Inches(3), Inches(6.0), Inches(7.3), Inches(0.01), theme["divider_color"], alpha=20)

    # Closing title - large and centered
    _add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
                  closing.get("title", "Thank You"), 48, theme["title_color"],
                  bold=True, alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    # Decorative divider line
    _add_shape_bg(slide, Inches(5.2), Inches(3.6), Inches(2.9), Inches(0.05), theme["accent"])

    # Closing bullets
    bullets = closing.get("bullets", [])
    y_start = 4.1
    for bullet in bullets:
        _add_text_box(slide, Inches(1), Inches(y_start), Inches(11.3), Inches(0.5),
                      bullet, 18, theme["subtitle_color"],
                      alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")
        y_start += 0.5

    # "Generated with SlideAI" branding
    _add_text_box(slide, Inches(4), Inches(6.8), Inches(5.3), Inches(0.4),
                  "Generated with SlideAI", 9, theme["slide_num_color"],
                  alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

    # Speaker notes
    notes = closing.get("speaker_notes", "")
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
        for paragraph in notes_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

    # Page number
    if total_slides > 0:
        _add_page_number(slide, total_slides, total_slides, theme)

    # Bottom accent bar
    _add_shape_bg(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), theme["bg_secondary"])


def generate_pptx(content: dict, template: str = "business") -> str:
    """
    Generate a .pptx file from structured content.

    Args:
        content: Dict with title, subtitle, author, slides, closing
        template: Theme name ('business', 'creative', 'minimal')

    Returns:
        Tuple of (file_id, filepath)
    """
    theme = THEMES.get(template, THEMES["business"])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create title slide
    _create_title_slide(prs, content, theme, template)

    # Create content slides
    slides = content.get("slides", [])
    total_slides = len(slides)
    for i, slide_data in enumerate(slides):
        _create_content_slide(prs, slide_data, i + 1, total_slides, theme, template)

    # Create closing slide
    closing = content.get("closing", {"title": "Thank You", "bullets": [], "speaker_notes": ""})
    _create_closing_slide(prs, closing, theme, template, total_slides=total_slides + 2)

    # Save file
    output_dir = os.path.join(os.path.dirname(__file__), "generated")
    os.makedirs(output_dir, exist_ok=True)

    file_id = str(uuid.uuid4())
    filename = f"{file_id}.pptx"
    filepath = os.path.join(output_dir, filename)

    prs.save(filepath)

    return file_id, filepath


def get_available_templates() -> list:
    """Return list of available templates with metadata."""
    return [
        {
            "id": "business",
            "name": "Business Professional",
            "description": "Clean, corporate look with navy blue tones. Ideal for boardroom presentations, investor pitches, and quarterly reviews.",
            "colors": ["#0F2B46", "#0096D6", "#00C9A7"]
        },
        {
            "id": "creative",
            "name": "Creative Modern",
            "description": "Bold gradients with purple and pink accents. Perfect for startups, marketing decks, and creative pitches.",
            "colors": ["#1A0A2E", "#A855F7", "#F472B6"]
        },
        {
            "id": "minimal",
            "name": "Minimal Clean",
            "description": "Elegant black and white with subtle grays. Best for academic presentations, technical talks, and formal reports.",
            "colors": ["#FAFAFA", "#181818", "#666666"]
        }
    ]
